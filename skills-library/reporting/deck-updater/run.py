"""deck-updater: update a PowerPoint deck with the latest numbers via {{KEY}} token replacement. Deterministic."""
import argparse
import json
import os
import sys

# --- locate the shared library (_shared/) whether run from its canonical path,
# --- a system's symlinked .claude/skills, or a standalone bundle -------------
_here = os.path.realpath(__file__)
_root = os.environ.get("IM_LIB_ROOT", "")
if not _root:
    _d = os.path.dirname(_here)
    while _d != os.path.dirname(_d):
        if os.path.isdir(os.path.join(_d, "_shared", "data-fetch")):
            _root = _d
            break
        _d = os.path.dirname(_d)
for _p in ("data-fetch", "router", "web-search"):
    _cand = os.path.join(_root, "_shared", _p)
    if os.path.isdir(_cand) and _cand not in sys.path:
        sys.path.insert(0, _cand)

from imdata import skillkit# noqa: E402

from pptx import Presentation  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402


def _parse_values(items):
    vals = {}
    for it in items or []:
        if "=" not in it:
            raise ValueError(f"--values entry '{it}' must be KEY=VALUE")
        k, v = it.split("=", 1)
        vals[k.strip()] = v
    return vals


def _build_base(values):
    """Create a simple 2-slide deck demonstrating the token mechanism."""
    prs = Presentation()
    # Slide 1: title with a {{TITLE}} token.
    title_layout = prs.slide_layouts[5]  # Title Only
    s1 = prs.slides.add_slide(title_layout)
    s1.shapes.title.text = "{{TITLE}}"
    # Slide 2: content listing the provided KEY: {{KEY}} pairs.
    blank = prs.slide_layouts[6]
    s2 = prs.slides.add_slide(blank)
    box = s2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(6))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for k in values:
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        para.text = f"{k}: {{{{{k}}}}}"
        para.font.size = Pt(20)
        first = False
    if first:  # no values supplied
        tf.paragraphs[0].text = "(no values provided)"
    return prs


def _replace_in_textframe(tf, values, counter):
    for para in tf.paragraphs:
        for run in para.runs:
            text = run.text
            for k, v in values.items():
                token = "{{" + k + "}}"
                if token in text:
                    counter[k] = counter.get(k, 0) + text.count(token)
                    text = text.replace(token, str(v))
            run.text = text


def main(args):
    values = {}
    if args.values_file:
        with open(args.values_file) as f:
            values.update({str(k): v for k, v in json.load(f).items()})
    values.update(_parse_values(args.values))

    if args.template:
        prs = Presentation(args.template)
        template_label = args.template
    else:
        prs = _build_base(values)
        template_label = None

    counter = {}
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                _replace_in_textframe(shape.text_frame, values, counter)

    out_dir = os.path.dirname(os.path.abspath(args.output))
    if out_dir and not os.path.isdir(out_dir):
        os.makedirs(out_dir, exist_ok=True)
    prs.save(args.output)

    replacements_made = sum(counter.values())
    keys_applied = sorted(counter.keys())
    slides = len(prs.slides._sldIdLst)

    if template_label:
        summary = (f"Updated deck from template '{template_label}': made "
                   f"{replacements_made} token replacement(s) across {slides} slide(s); "
                   f"saved to {args.output}.")
    else:
        summary = (f"Built a {slides}-slide deck from scratch and applied "
                   f"{replacements_made} token replacement(s); saved to {args.output}.")

    return {
        "template": template_label,
        "output": os.path.abspath(args.output),
        "replacements_made": replacements_made,
        "keys_applied": keys_applied,
        "slides": slides,
        "summary": summary,
    }


if __name__ == "__main__":
    p = argparse.ArgumentParser(description="Update a .pptx deck via {{KEY}} replacement.")
    p.add_argument("--template", default=None,
                   help="path to a .pptx; if omitted, build a 2-slide base deck")
    p.add_argument("--output", required=True, help="path to write the updated .pptx")
    p.add_argument("--values", nargs="*", default=None, help="repeatable KEY=VALUE")
    p.add_argument("--values-file", default=None, help="JSON file of {KEY: VALUE}")
    skillkit.run(main, p)
